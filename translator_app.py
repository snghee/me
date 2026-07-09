from __future__ import annotations

import os
import threading
import zipfile
from io import BytesIO
from pathlib import Path
from typing import Callable, Optional
from xml.etree import ElementTree as ET

try:
    import tkinter as tk
except ModuleNotFoundError:  # pragma: no cover - exercised in CLI-only environments
    tk = None

import requests
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

SUPPORTED_EXTENSIONS = {".txt", ".docx", ".pdf", ".pptx"}
API_URL = "https://libretranslate.de/translate"


class TranslationApp:
    def __init__(self, root) -> None:
        if tk is None:
            raise RuntimeError("Tkinter is not available in this environment.")

        self.root = root
        self.root.title("문서 번역 프로그램")
        self.root.geometry("760x460")
        self.root.configure(bg="#f4f7fb")

        self.source_path: Optional[Path] = None
        self.status_var = tk.StringVar(value="준비중")
        self.file_var = tk.StringVar(value="선택된 원본 파일 없음")
        self.output_var = tk.StringVar(value="")
        self.source_lang_var = tk.StringVar(value="auto")
        self.target_lang_var = tk.StringVar(value="ja")
        self.output_dir_var = tk.StringVar(value="")

        self._build_ui()

    def _build_ui(self) -> None:
        title = tk.Label(
            self.root,
            text="일본어 / 한국어 문서 번역기",
            font=("Malgun Gothic", 20, "bold"),
            bg="#f4f7fb",
            fg="#1f2937",
        )
        title.pack(pady=(24, 8))

        subtitle = tk.Label(
            self.root,
            text="파일 탐색기를 통해 원본 문서를 선택하고 번역을 시작하세요.",
            font=("Malgun Gothic", 12),
            bg="#f4f7fb",
            fg="#4b5563",
        )
        subtitle.pack(pady=(0, 20))

        card = tk.Frame(self.root, bg="white", bd=1, relief="solid")
        card.pack(fill="both", expand=True, padx=24, pady=(0, 24))

        tk.Label(card, text="진행 상태", font=("Malgun Gothic", 13, "bold"), bg="white", fg="#111827").pack(anchor="w", padx=20, pady=(18, 8))
        status_label = tk.Label(card, textvariable=self.status_var, font=("Malgun Gothic", 14), bg="white", fg="#2563eb")
        status_label.pack(anchor="w", padx=20, pady=(0, 12))

        tk.Label(card, text="원본 파일", font=("Malgun Gothic", 12, "bold"), bg="white", fg="#374151").pack(anchor="w", padx=20, pady=(0, 4))
        file_label = tk.Label(card, textvariable=self.file_var, font=("Malgun Gothic", 11), bg="white", fg="#4b5563", wraplength=680, justify="left")
        file_label.pack(anchor="w", padx=20, pady=(0, 8))

        tk.Label(card, text="결과 파일", font=("Malgun Gothic", 12, "bold"), bg="white", fg="#374151").pack(anchor="w", padx=20, pady=(0, 4))
        output_label = tk.Label(card, textvariable=self.output_var, font=("Malgun Gothic", 11), bg="white", fg="#059669", wraplength=680, justify="left")
        output_label.pack(anchor="w", padx=20, pady=(0, 18))

        option_frame = tk.Frame(card, bg="white")
        option_frame.pack(fill="x", padx=20, pady=(0, 10))
        tk.Label(option_frame, text="원본 언어", bg="white", font=("Malgun Gothic", 11, "bold"), fg="#374151").pack(side="left")
        source_menu = tk.OptionMenu(option_frame, self.source_lang_var, "auto", "ko", "ja", "en")
        source_menu.config(width=10, font=("Malgun Gothic", 10))
        source_menu.pack(side="left", padx=(8, 0))

        tk.Label(option_frame, text="대상 언어", bg="white", font=("Malgun Gothic", 11, "bold"), fg="#374151").pack(side="left", padx=(16, 8))
        target_menu = tk.OptionMenu(option_frame, self.target_lang_var, "ja", "ko", "en")
        target_menu.config(width=10, font=("Malgun Gothic", 10))
        target_menu.pack(side="left")

        tk.Label(option_frame, text="저장 폴더", bg="white", font=("Malgun Gothic", 11, "bold"), fg="#374151").pack(side="left", padx=(16, 8))
        folder_entry = tk.Entry(option_frame, textvariable=self.output_dir_var, width=32)
        folder_entry.pack(side="left")
        tk.Button(option_frame, text="찾기", command=self.select_output_dir, width=8, bg="#e5e7eb", relief="flat", font=("Malgun Gothic", 10)).pack(side="left", padx=(8, 0))

        button_frame = tk.Frame(card, bg="white")
        button_frame.pack(fill="x", padx=20, pady=(0, 20))

        tk.Button(button_frame, text="원본 파일 선택", command=self.select_file, width=18, height=1, bg="#2563eb", fg="white", relief="flat", font=("Malgun Gothic", 11, "bold")).pack(side="left", padx=(0, 10))
        tk.Button(button_frame, text="번역 시작", command=self.start_translation, width=18, height=1, bg="#059669", fg="white", relief="flat", font=("Malgun Gothic", 11, "bold")).pack(side="left")

    def select_file(self) -> None:
        file_path = choose_source_file()
        if not file_path:
            return

        self.source_path = file_path
        self.file_var.set(str(self.source_path))
        self.output_var.set("")
        self.status_var.set("원본선택")

    def select_output_dir(self) -> None:
        from tkinter import filedialog

        folder_path = filedialog.askdirectory(title="저장 폴더 선택")
        if folder_path:
            self.output_dir_var.set(folder_path)

    def start_translation(self) -> None:
        if not self.source_path:
            self.status_var.set("원본선택")
            return

        self.status_var.set("번역중")
        self.output_var.set("번역을 진행 중입니다...")
        thread = threading.Thread(target=self._run_translation, daemon=True)
        thread.start()

    def _run_translation(self) -> None:
        try:
            output_dir = Path(self.output_dir_var.get()).expanduser() if self.output_dir_var.get() else None
            output_path = translate_document(
                self.source_path,
                source_lang=self.source_lang_var.get(),
                target_lang=self.target_lang_var.get(),
                output_dir=output_dir,
            )
            self.root.after(0, lambda: self._finish_translation(output_path))
        except Exception as exc:  # pragma: no cover - GUI error path
            self.root.after(0, lambda: self._fail_translation(str(exc)))

    def _finish_translation(self, output_path: Path) -> None:
        self.status_var.set("완료")
        self.output_var.set(str(output_path))

    def _fail_translation(self, message: str) -> None:
        self.status_var.set("준비중")
        self.output_var.set(f"번역 실패: {message}")


def choose_source_file() -> Optional[Path]:
    if tk is None:
        return None

    from tkinter import filedialog

    file_path = filedialog.askopenfilename(
        title="원본 문서 선택",
        filetypes=[
            ("텍스트 문서", "*.txt"),
            ("워드 문서", "*.docx"),
            ("PDF 문서", "*.pdf"),
            ("파워포인트", "*.pptx"),
        ],
    )
    if not file_path:
        return None
    return Path(file_path)


def build_output_path(source_path: Path, target_lang: str, output_dir: Optional[Path] = None) -> Path:
    suffix = get_output_suffix(target_lang)
    output_path = source_path.with_name(f"{source_path.stem}_{suffix}{source_path.suffix}")
    if output_dir is not None:
        output_dir.mkdir(parents=True, exist_ok=True)
        return output_dir / output_path.name
    return output_path


def get_output_suffix(target_lang: str) -> str:
    if target_lang == "ja":
        return "일어"
    if target_lang == "ko":
        return "한국어"
    if target_lang == "en":
        return "영어"
    return "번역"


def split_text_chunks(text: str, max_chars: int = 3000) -> list[str]:
    if len(text) <= max_chars:
        return [text]

    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + max_chars)
        chunk = text[start:end]
        if end < len(text):
            last_break = max(chunk.rfind("\n"), chunk.rfind("."), chunk.rfind("!"), chunk.rfind("?"))
            if last_break > 0:
                end = start + last_break + 1
                chunk = text[start:end]
        chunks.append(chunk)
        start = end
    return chunks


def translate_text(text: str, source_lang: str = "auto", target_lang: str = "ja", translator: Optional[Callable[..., str]] = None) -> str:
    if not text or not text.strip():
        return text

    if translator:
        try:
            return translator(text, source_lang, target_lang)
        except TypeError:
            return translator(text)

    chunks = split_text_chunks(text)
    translated_parts: list[str] = []
    for chunk in chunks:
        translated_parts.append(_translate_with_libretranslate(chunk, source_lang, target_lang))
    return "\n".join(translated_parts)


def _translate_with_libretranslate(text: str, source_lang: str, target_lang: str) -> str:
    try:
        response = requests.post(
            API_URL,
            json={"q": text, "source": source_lang, "target": target_lang, "format": "text"},
            timeout=30,
        )
        response.raise_for_status()
        payload = response.json()
        if isinstance(payload, dict) and "translatedText" in payload:
            return payload["translatedText"]
        if isinstance(payload, dict) and isinstance(payload.get("data"), dict):
            return payload["data"].get("translatedText", text)
    except Exception:
        pass

    return text


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".txt":
        return path.read_text(encoding="utf-8")
    if suffix == ".docx":
        document = Document(path)
        paragraphs: list[str] = []
        for paragraph in document.paragraphs:
            if paragraph.text.strip():
                paragraphs.append(paragraph.text)
        return "\n".join(paragraphs)
    if suffix == ".pptx":
        presentation = Presentation(path)
        slides: list[str] = []
        for slide in presentation.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text.strip():
                    texts.append(shape.text)
            if texts:
                slides.append("\n".join(texts))
        return "\n\n".join(slides)
    if suffix == ".pdf":
        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    raise ValueError(f"지원되지 않는 파일 형식입니다: {path.suffix}")


def translate_text_file(source_path: Path, source_lang: str = "auto", target_lang: str = "ja", translator: Optional[Callable[[str, str, str], str]] = None, output_dir: Optional[Path] = None) -> Path:
    text = extract_text(source_path)
    translated_text = translate_text(text, source_lang=source_lang, target_lang=target_lang, translator=translator)
    output_path = build_output_path(source_path, target_lang, output_dir=output_dir)
    output_path.write_text(translated_text, encoding="utf-8")
    return output_path


def translate_document(source_path: Path, source_lang: str = "auto", target_lang: str = "ja", output_dir: Optional[Path] = None) -> Path:
    if source_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
        raise ValueError("지원 가능한 형식은 TXT, DOCX, PDF, PPTX 입니다.")

    if source_path.suffix.lower() == ".txt":
        return translate_text_file(source_path, source_lang, target_lang, output_dir=output_dir)

    if source_path.suffix.lower() == ".docx":
        output_path = build_output_path(source_path, target_lang, output_dir=output_dir)
        translate_docx_with_images(source_path, output_path, source_lang=source_lang, target_lang=target_lang)
        return output_path

    if source_path.suffix.lower() == ".pptx":
        source_presentation = Presentation(source_path)
        translated_presentation = Presentation()
        blank_layout = translated_presentation.slide_layouts[6]

        for slide in source_presentation.slides:
            new_slide = translated_presentation.slides.add_slide(blank_layout)
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text.strip():
                    textbox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                    text_frame = textbox.text_frame
                    text_frame.clear()
                    translated_text = translate_text(shape.text, source_lang=source_lang, target_lang=target_lang)
                    paragraph = text_frame.paragraphs[0]
                    paragraph.text = translated_text
                elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    image_bytes = shape.image.blob
                    new_slide.shapes.add_picture(BytesIO(image_bytes), shape.left, shape.top, shape.width, shape.height)

        output_path = build_output_path(source_path, target_lang, output_dir=output_dir)
        translated_presentation.save(output_path)
        return output_path

    if source_path.suffix.lower() == ".pdf":
        text = extract_text(source_path)
        translated_text = translate_text(text, source_lang=source_lang, target_lang=target_lang)
        output_path = build_output_path(source_path, target_lang, output_dir=output_dir).with_suffix('.pdf')
        write_pdf(output_path, translated_text)
        return output_path

    raise ValueError("지원되지 않는 파일 형식입니다.")


def translate_docx_with_images(source_path: Path, output_path: Path, source_lang: str = "auto", target_lang: str = "ja") -> None:
    with zipfile.ZipFile(source_path, "r") as src_zip:
        with zipfile.ZipFile(output_path, "w") as dst_zip:
            for item in src_zip.infolist():
                if item.filename == "word/document.xml":
                    xml_bytes = src_zip.read(item.filename)
                    root = ET.fromstring(xml_bytes)
                    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
                    for node in root.findall('.//w:t', ns):
                        if node.text and node.text.strip():
                            node.text = translate_text(node.text, source_lang=source_lang, target_lang=target_lang)
                    translated_xml = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                    dst_zip.writestr(item, translated_xml)
                else:
                    dst_zip.writestr(item, src_zip.read(item.filename))


def write_pdf(output_path: Path, text: str) -> None:
    doc = SimpleDocTemplate(str(output_path), pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    for paragraph in text.splitlines():
        if paragraph.strip():
            story.append(Paragraph(paragraph, styles["BodyText"]))
            story.append(Spacer(1, 6))
    doc.build(story)


def main() -> None:
    import argparse

    parser = argparse.ArgumentParser(description="문서 번역 프로그램")
    parser.add_argument("source", nargs="?", help="번역할 원본 파일 경로")
    parser.add_argument("--source-lang", default="auto", help="원본 언어")
    parser.add_argument("--target-lang", default="ja", help="대상 언어")
    parser.add_argument("--output-dir", default="", help="번역 결과 저장 폴더")
    args = parser.parse_args()

    if args.source:
        source_path = Path(args.source).expanduser().resolve()
        output_dir = Path(args.output_dir).expanduser().resolve() if args.output_dir else None
        output_path = translate_document(source_path, source_lang=args.source_lang, target_lang=args.target_lang, output_dir=output_dir)
        print(output_path)
        return

    if tk is None:
        print("GUI 환경이 없어 CLI 모드로 실행되었습니다. 파일 경로를 인자로 전달해 주세요.")
        return

    try:
        root = tk.Tk()
        TranslationApp(root)
        root.mainloop()
    except Exception:
        print("GUI 환경이 없어 CLI 모드로 실행되었습니다. 파일 경로를 인자로 전달해 주세요.")


if __name__ == "__main__":
    main()
